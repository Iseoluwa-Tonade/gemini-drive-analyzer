import os
import streamlit as st
import json
import google.generativeai as genai
from google_auth_oauthlib.flow import Flow
from googleapiclient.discovery import build
import io
import fitz  # PyMuPDF
import docx
import pptx
from PIL import Image
from google.oauth2.credentials import Credentials

# --- CONFIG ---
SCOPES = ['https://www.googleapis.com/auth/drive.readonly']
REDIRECT_URI = "https://zw2bm6uwryon2f5pnfsauk.streamlit.app"

# --- HELPER FUNCTION ---
def get_file_content(drive_service, file_info):
    """
    Downloads and extracts content from a Google Drive file.
    Handles text, pdf, docx, pptx, and images.
    """
    file_id = file_info['id']
    mime_type = file_info.get('mimeType', '')

    try:
        request = drive_service.files().get_media(fileId=file_id)
        file_bytes = request.execute()

        if mime_type.startswith("text/"):
            return 'text', file_bytes.decode("utf-8", errors="ignore")

        elif mime_type == "application/pdf":
            pdf_doc = fitz.open(stream=file_bytes, filetype="pdf")
            text = ""
            for page in pdf_doc:
                text += page.get_text()
            return 'text', text

        elif mime_type in [
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword"
        ]:
            doc = docx.Document(io.BytesIO(file_bytes))
            text = "\n".join([p.text for p in doc.paragraphs])
            return 'text', text

        elif mime_type in [
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ]:
            prs = pptx.Presentation(io.BytesIO(file_bytes))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return 'text', text

        elif mime_type.startswith("image/"):
            return 'image', file_bytes

        else:
            return 'unsupported', f"Unsupported MIME type: {mime_type}"

    except Exception as e:
        return 'unsupported', f"Error reading file: {e}"

# --- GEMINI CALL ---
def get_gemini_response(api_key, prompt_parts):
    """
    Sends a multimodal prompt to the Gemini API.
    prompt_parts = list of strings and/or images.
    """
    try:
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel(model_name="gemini-1.5-pro-latest")
        response = model.generate_content(prompt_parts)
        return response.text
    except Exception as e:
        return f"ERROR: Could not generate response from Gemini. Details: {str(e)}"

# --- MAIN APP ---
st.set_page_config(page_title="Drive + Gemini Analyzer", page_icon="✨")

st.title("✨ Analyze Google Drive Files with Gemini ✨")

# --- AUTH ---
if "credentials" not in st.session_state:
    st.session_state.credentials = None

if not st.session_state.credentials:
    try:
        flow = Flow.from_client_config(
            client_config=st.secrets["google_credentials"],
            scopes=SCOPES,
            redirect_uri=REDIRECT_URI
        )
        auth_url, _ = flow.authorization_url(prompt='consent')
        st.markdown(f"[Authorize Google Drive]({auth_url})")

        # ✅ NEW QUERY PARAMS LOGIC
        query_params = st.query_params
        if "code" in query_params:
            code = query_params["code"]
            flow.fetch_token(code=code)
            creds = flow.credentials
            st.session_state.credentials = creds.to_json()
            st.rerun()

    except KeyError:
        st.error('The "google_credentials" secret is missing. Please add it to your Streamlit secrets.')
    except Exception as e:
        st.error("Could not load Google credentials from secrets. Ensure they are correctly configured.")
        st.error(f"Specific error: {e}")
else:
    creds = Credentials.from_authorized_user_info(
        json.loads(st.session_state.credentials)
    )
    drive_service = build('drive', 'v3', credentials=creds)

    # --- FETCH FILES ---
    results = drive_service.files().list(pageSize=50, fields="files(id, name, mimeType)").execute()
    items = results.get('files', [])

    if not items:
        st.info("No files found in your Google Drive.")
    else:
        file_options = {item["name"]: item for item in items}

        selected_files_display = st.multiselect(
            "Choose files to analyze:",
            options=list(file_options.keys())
        )

        user_prompt = st.text_area("What would you like to know about these files?", height=100)

        if st.button("✨ Analyze Files with Gemini", disabled=(not selected_files_display or not user_prompt)):
            prompt_parts = []
            with st.status("Processing files...", expanded=True) as status:
                for file_display in selected_files_display:
                    file_info = file_options[file_display]
                    status.update(label=f"Processing: {file_display}...")

                    content_type, content = get_file_content(drive_service, file_info)

                    if content_type == 'text':
                        prompt_parts.append(f"\n--- DOCUMENT: {file_info['name']} ---\n{content}")
                    elif content_type == 'image':
                        img = Image.open(io.BytesIO(content))
                        prompt_parts.append(img)
                    else:
                        st.warning(f"Skipping unsupported file: {file_info['name']} ({content})")

                status.update(label="All files processed!", state="complete")

            if prompt_parts:
                prompt_parts.insert(0, user_prompt)

                st.info(f"Sending {len(selected_files_display)} file(s) to Gemini for analysis...")
                with st.spinner("🤖 Gemini is thinking..."):
                    api_key = st.secrets["GOOGLE_API_KEY"]
                    gemini_response = get_gemini_response(api_key, prompt_parts)
                    st.subheader("Gemini's Response:")
                    st.write(gemini_response)
